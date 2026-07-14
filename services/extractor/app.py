"""Microserviço de extração de texto do DMDoc (Python).

Modos de operação:
  - Fila Redis (primário): consumer loop BRPOP em 'extract:requests'; baixa o
    arquivo do S3 diretamente; publica resultado em 'extract:result:{requestId}'.
    Elimina timeout HTTP — o worker aguarda via BLPOP sem limite de tempo.
  - Endpoint HTTP /extract (legado): mantido para compatibilidade e testes manuais.
  - Endpoint HTTP /convert/pdf: conversão Office→PDF para preview, não afetado.

Dispatch por formato:
- PDF   → PyMuPDF (texto nativo); páginas escaneadas (sem texto) caem para OCR
- JPG/PNG/WebP → OCR (Tesseract) com remoção de borda branca + auto-rotação via OSD
- DOCX  → python-docx (parágrafos + tabelas) + OCR de imagens embutidas
- XLSX  → openpyxl (células por planilha)
- PPTX  → python-pptx (texto dos slides)

OCR: Tesseract com LSTM engine (--oem 1). Auto-rotação via Tesseract OSD (--psm 0)
em vez de EasyOCR/PyTorch — funciona sem AVX2/FMA e processa 1-3s/página em CPUs
antigas (vs. 20-90s do EasyOCR em Westmere/Ivy Bridge sem AVX2).
"""

import asyncio
import datetime
import io
import json
import logging
import os
import subprocess
import tempfile
import threading
import zipfile
from zoneinfo import ZoneInfo

import boto3
import redis as redis_lib

import cv2
import pytesseract
import fitz  # PyMuPDF
import numpy as np
from PIL import Image

_SAO_PAULO_TZ = ZoneInfo("America/Sao_Paulo")

# Atributos padrão de um LogRecord — tudo que NÃO estiver aqui é tratado como
# campo de contexto (tenantId, documentId, traceId, ...) passado via `extra=`.
# `color_message` é ruído do uvicorn (mensagem com códigos ANSI) — descartado.
_LOG_RECORD_BUILTINS = frozenset(vars(logging.makeLogRecord({}))) | {
    "message",
    "asctime",
    "taskName",
    "color_message",
}


class JsonFormatter(logging.Formatter):
    """Formata logs como JSON no padrão único do DMDoc (ver packages/logger).

    Cada linha contém:
      - level:   label minúsculo (info, error, ...)
      - time:    yyyy-mm-dd hh:mm:ss no fuso America/Sao_Paulo
      - service: "extractor"
      - msg:     mensagem já interpolada
      - campos de contexto extras (traceId, tenantId, documentId, ...)
    """

    def format(self, record: logging.LogRecord) -> str:
        payload: dict[str, object] = {
            "level": record.levelname.lower(),
            "time": datetime.datetime.fromtimestamp(record.created, _SAO_PAULO_TZ).strftime(
                "%Y-%m-%d %H:%M:%S"
            ),
            "service": "extractor",
            "msg": record.getMessage(),
        }
        for key, value in record.__dict__.items():
            if key not in _LOG_RECORD_BUILTINS and not key.startswith("_"):
                payload[key] = value
        if record.exc_info:
            payload["err"] = self.formatException(record.exc_info)
        return json.dumps(payload, ensure_ascii=False)


def _install_json_handler() -> logging.Logger:
    """Instala um handler JSON no logger raiz e roteia uvicorn/FastAPI por ele,
    para que toda a saída do serviço siga o mesmo padrão de log."""
    handler = logging.StreamHandler()
    handler.setFormatter(JsonFormatter())

    root = logging.getLogger()
    root.handlers = [handler]
    root.setLevel(logging.INFO)

    # uvicorn configura seus loggers antes de importar este módulo; aqui os
    # esvaziamos e deixamos propagar para o handler JSON do root.
    for name in ("uvicorn", "uvicorn.error", "uvicorn.access"):
        uv = logging.getLogger(name)
        uv.handlers = []
        uv.propagate = True

    app_logger = logging.getLogger("dmdoc.extractor")
    app_logger.setLevel(logging.INFO)
    app_logger.propagate = True
    return app_logger


logger = _install_json_handler()

LANGS = os.environ.get("OCR_LANGS", "pt").split(",")

# Mapeia códigos EasyOCR/ISO → Tesseract (mantém compatibilidade com OCR_LANGS existente)
_LANG_MAP = {"pt": "por", "en": "eng", "es": "spa", "fr": "fra", "de": "deu"}
OCR_LANG_TESS = "+".join(_LANG_MAP.get(l.strip(), l.strip()) for l in LANGS)

# Timeout máximo (segundos) para o endpoint /extract legado (HTTP direto).
# O consumer Redis não usa este timeout — processa sem limite de tempo.
REQUEST_TIMEOUT_S = int(os.environ.get("EXTRACT_TIMEOUT_S", "120"))

# Configuração do consumer Redis
REDIS_URL = os.environ.get("REDIS_URL", "redis://localhost:6379")
S3_ENDPOINT = os.environ.get("S3_ENDPOINT", "")
S3_BUCKET = os.environ.get("S3_BUCKET", "dmdoc-documents")
AWS_REGION = os.environ.get("AWS_REGION", "us-east-1")
AWS_ACCESS_KEY_ID = os.environ.get("AWS_ACCESS_KEY_ID", "")
AWS_SECRET_ACCESS_KEY = os.environ.get("AWS_SECRET_ACCESS_KEY", "")
# TTL do result key no Redis — cleanup automático se o worker morrer antes de ler.
RESULT_KEY_TTL_SECS = 3600
EXTRACT_REQUEST_QUEUE = "extract:requests"
EXTRACT_RESULT_PREFIX = "extract:result:"

# Mapeamento content-type → extensão de arquivo para conversão Office→PDF.
OFFICE_EXTENSIONS: dict[str, str] = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.ms-powerpoint": ".ppt",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/msword": ".doc",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.ms-excel": ".xls",
    "application/vnd.oasis.opendocument.presentation": ".odp",
    "application/vnd.oasis.opendocument.text": ".odt",
    "text/rtf": ".rtf",
    "application/rtf": ".rtf",
}

# Resolução máxima enviada ao OCR (maior lado, px).
MAX_DIM = 1600
# Limiar de grayscale: pixel < CONTENT_THRESHOLD conta como conteúdo (não-branco).
CONTENT_THRESHOLD = 220
# Mínimo de caracteres alfanuméricos numa página de PDF para considerá-la "nativa"
# (com camada de texto). Abaixo disso, a página é tratada como escaneada → OCR.
MIN_PAGE_ALNUM = 20

# OCR complementar em páginas de PDF híbridas (texto nativo + imagem embutida com
# dados, ex.: boleto/nota fiscal renderizado como imagem dentro de um PDF que já
# tem cabeçalho/rodapé/título em texto nativo). Uma página que passa no critério
# de texto nativo ainda dispara OCR complementar quando:
#   - as imagens embutidas cobrem >= PDF_IMG_COVERAGE_MIN da área da página; E
#   - o texto nativo cobre <= PDF_NATIVE_TEXT_COVERAGE_MAX da área da página.
# Assim evita-se OCR desnecessário em logos/ícones pequenos (baixa cobertura de
# imagem) e em páginas densas de texto que só têm um gráfico ilustrativo (alta
# cobertura de texto nativo), sem baixar o MIN_PAGE_ALNUM (que reintroduziria
# falsos positivos em páginas já bem cobertas por texto).
PDF_IMG_COVERAGE_MIN = 0.15
PDF_NATIVE_TEXT_COVERAGE_MAX = 0.60
# Resolução (células no maior lado) do grid booleano usado para medir cobertura
# de imagem/texto sem dupla contagem de retângulos sobrepostos.
PDF_COVERAGE_GRID = 120

# MIMEs e extensões de imagem suportados pelo cv2.imdecode (libwebp embutida).
_IMAGE_MIMES = frozenset({
    "image/jpeg", "image/png", "image/webp",
    "image/gif", "image/bmp", "image/tiff", "image/tif",
})
_IMAGE_EXTS = (".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tiff", ".tif")

from docx import Document as DocxDocument
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import Response
from openpyxl import load_workbook
from pptx import Presentation

app = FastAPI(title="DMDoc Extractor", version="1.0.0")


# ──────────────────────────── Redis consumer ────────────────────────────


def _make_s3_client():
    kwargs = dict(
        region_name=AWS_REGION,
        aws_access_key_id=AWS_ACCESS_KEY_ID or None,
        aws_secret_access_key=AWS_SECRET_ACCESS_KEY or None,
    )
    if S3_ENDPOINT:
        kwargs["endpoint_url"] = S3_ENDPOINT
    return boto3.client("s3", **kwargs)


def _extraction_consumer_loop() -> None:
    """Consumer Redis que processa pedidos de extração sem timeout HTTP.

    Roda em thread background (daemon=True) iniciada no startup do FastAPI.
    BRPOP bloqueia indefinidamente até chegar um pedido; processa um job por vez,
    serializando OCR naturalmente sem necessidade de asyncio.Semaphore.
    Em erro de conexão: dorme 5s e reconecta.
    Em erro de extração: publica {error} no result key para o worker não ficar preso.
    """
    import time

    rconn = redis_lib.from_url(REDIS_URL, decode_responses=True)
    s3 = _make_s3_client()
    logger.info("extraction consumer loop iniciado (fila=%s)", EXTRACT_REQUEST_QUEUE)

    while True:
        try:
            item = rconn.brpop(EXTRACT_REQUEST_QUEUE, timeout=0)
            if item is None:
                continue

            _, raw = item
            try:
                payload = json.loads(raw)
            except json.JSONDecodeError as exc:
                logger.error("payload JSON inválido na fila: %.200s — %s", raw, exc)
                continue

            request_id = payload.get("requestId", "")
            s3_key = payload.get("s3Key", "")
            s3_bucket = payload.get("s3Bucket", S3_BUCKET)
            mime = payload.get("mimeType", "")
            result_key = f"{EXTRACT_RESULT_PREFIX}{request_id}"

            logger.info(
                "extraindo s3Key=%s mime=%s",
                s3_key, mime,
                extra={"traceId": request_id, "s3Key": s3_key, "mimeType": mime},
            )

            try:
                s3_resp = s3.get_object(Bucket=s3_bucket, Key=s3_key)
                data = s3_resp["Body"].read()
                name = s3_key.split("/")[-1].lower()
                result = _dispatch_extract(data, mime, name)
                result["engine"] = "python"
                out = json.dumps(result)
            except Exception as exc:  # noqa: BLE001
                logger.exception(
                    "falha ao processar s3Key=%s: %s",
                    s3_key, exc,
                    extra={"traceId": request_id, "s3Key": s3_key},
                )
                out = json.dumps({"error": str(exc), "text": "", "pageCount": 1, "ocrPages": []})

            pipe = rconn.pipeline()
            pipe.rpush(result_key, out)
            pipe.expire(result_key, RESULT_KEY_TTL_SECS)
            pipe.execute()
            logger.info(
                "resultado publicado", extra={"traceId": request_id}
            )

        except redis_lib.exceptions.ConnectionError as exc:
            logger.error(
                "Redis connection error no consumer loop: %s — reconectando em 5s", exc
            )
            time.sleep(5)
            try:
                rconn = redis_lib.from_url(REDIS_URL, decode_responses=True)
                s3 = _make_s3_client()
            except Exception:  # noqa: BLE001
                logger.exception("falha ao reconectar Redis/S3")
        except Exception as exc:  # noqa: BLE001
            logger.exception("erro inesperado no consumer loop: %s", exc)
            time.sleep(1)


@app.on_event("startup")
def _start_consumer() -> None:
    t = threading.Thread(
        target=_extraction_consumer_loop,
        daemon=True,
        name="extraction-consumer",
    )
    t.start()
    logger.info("extraction consumer thread iniciado")


# ──────────────────────────── OCR (Tesseract) ────────────────────────────


def _crop_white(img: np.ndarray) -> np.ndarray:
    """Recorta a borda branca (folha em volta do documento) via bbox de pixels
    não-brancos, com mediana para ignorar sujeira isolada no branco."""
    h, w = img.shape[:2]
    if w == 0 or h == 0:
        return img
    sw = 200
    sh = max(1, round(h / w * sw))
    thumb = cv2.cvtColor(cv2.resize(img, (sw, sh)), cv2.COLOR_BGR2GRAY)
    thumb = cv2.medianBlur(thumb, 5)
    ys, xs = np.where(thumb < CONTENT_THRESHOLD)
    if xs.size < 50:
        return img
    minx, maxx, miny, maxy = int(xs.min()), int(xs.max()), int(ys.min()), int(ys.max())
    sx, sy, pad = w / sw, h / sh, 4
    left = max(0, int((minx - pad) * sx))
    top = max(0, int((miny - pad) * sy))
    right = min(w, int((maxx + pad) * sx))
    bot = min(h, int((maxy + pad) * sy))
    if right <= left or bot <= top:
        return img
    if (right - left) * (bot - top) >= w * h * 0.95:
        return img
    return img[top:bot, left:right]


def _resize(img: np.ndarray, max_dim: int) -> np.ndarray:
    """Reduz a imagem para que o maior lado seja no máximo `max_dim` (nunca amplia)."""
    h, w = img.shape[:2]
    scale = min(max_dim / max(h, w), 1.0)
    if scale < 1.0:
        img = cv2.resize(img, (int(w * scale), int(h * scale)), interpolation=cv2.INTER_AREA)
    return img


def _to_pil(img: np.ndarray) -> Image.Image:
    """Converte imagem BGR (OpenCV) para PIL RGB (Tesseract/Pillow)."""
    return Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))


def _auto_rotate(pil_img: Image.Image) -> Image.Image:
    """Detecta orientação via Tesseract OSD (--psm 0) e corrige. Falha silenciosamente."""
    try:
        osd = pytesseract.image_to_osd(
            pil_img, output_type=pytesseract.Output.DICT, config="--psm 0"
        )
        angle = osd.get("rotate", 0)
        if angle:
            pil_img = pil_img.rotate(angle, expand=True)
    except Exception:
        pass
    return pil_img


def ocr_image(img: np.ndarray) -> str:
    """OCR via Tesseract LSTM:
    1. Recorta borda branca;
    2. Reduz para MAX_DIM;
    3. Detecta e corrige orientação via Tesseract OSD (barato, sem rede neural pesada);
    4. Roda reconhecimento com --oem 1 (LSTM) --psm 3 (auto page segmentation).

    Toda a etapa é protegida por try/except: uma imagem patológica retorna
    texto vazio + log, em vez de derrubar o serviço."""
    try:
        cropped = _crop_white(img)
        full = _resize(cropped, MAX_DIM)
        pil_img = _auto_rotate(_to_pil(full))
        return pytesseract.image_to_string(
            pil_img, lang=OCR_LANG_TESS, config="--oem 1 --psm 3"
        ).strip()
    except Exception:  # noqa: BLE001
        logger.exception("ocr_image falhou; retornando texto vazio para este item")
        return ""


# ──────────────────────────── extração por formato ────────────────────────────


def _libreoffice_to_pdf(data: bytes, ext: str) -> bytes:
    """Converte documento Office/RTF para PDF via LibreOffice headless."""
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, f"input{ext}")
        with open(input_path, "wb") as f:
            f.write(data)
        try:
            result = subprocess.run(
                ["libreoffice", "--headless", "--norestore",
                 "--convert-to", "pdf", "--outdir", tmpdir, input_path],
                capture_output=True, timeout=60,
            )
        except subprocess.TimeoutExpired:
            return b""
        if result.returncode != 0:
            return b""
        pdf_name = os.path.basename(input_path).rsplit(".", 1)[0] + ".pdf"
        pdf_path = os.path.join(tmpdir, pdf_name)
        if not os.path.exists(pdf_path):
            return b""
        with open(pdf_path, "rb") as f:
            return f.read()


def extract_txt(data: bytes) -> dict:
    text = data.decode("utf-8", errors="replace")
    return {"text": text, "pageCount": 1, "ocrPages": []}


def _ocr_page_pixmap(page) -> str:
    """Renderiza a página inteira (200 dpi) e roda OCR sobre a imagem composta.

    Renderizar a página composta — em vez de extrair cada imagem embutida via
    page.get_images() — é mais robusto para boletos/formulários: as bandas de
    imagem costumam vir empilhadas, com máscaras/canais separados e escalas
    próprias; extraí-las cruas produz camadas parciais e texto ilegível. A
    página renderizada reflete exatamente o que o usuário vê e reaproveita o
    mesmo caminho de OCR já usado para páginas escaneadas."""
    pix = page.get_pixmap(dpi=200)
    img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
    if pix.n == 4:
        img = cv2.cvtColor(img, cv2.COLOR_RGBA2BGR)
    elif pix.n == 3:
        img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)
    else:
        img = cv2.cvtColor(img, cv2.COLOR_GRAY2BGR)
    return ocr_image(img)


def _page_coverage(page) -> tuple[float, float]:
    """Retorna (fração da página coberta por imagens embutidas, fração coberta
    por texto nativo).

    Usa um grid booleano de baixa resolução para calcular a área de união dos
    retângulos, evitando dupla contagem de imagens sobrepostas (bandas empilhadas
    de boleto) e de blocos de texto que se tocam."""
    rect = page.rect
    if rect.width <= 0 or rect.height <= 0:
        return 0.0, 0.0
    gw = PDF_COVERAGE_GRID
    gh = max(1, round(rect.height / rect.width * gw))
    sx = gw / rect.width
    sy = gh / rect.height

    def _covered_frac(rects: list) -> float:
        if not rects:
            return 0.0
        mask = np.zeros((gh, gw), dtype=bool)
        for r in rects:
            x0 = max(0, int(r.x0 * sx))
            x1 = min(gw, int(round(r.x1 * sx)))
            y0 = max(0, int(r.y0 * sy))
            y1 = min(gh, int(round(r.y1 * sy)))
            if x1 > x0 and y1 > y0:
                mask[y0:y1, x0:x1] = True
        return float(mask.sum()) / float(gw * gh)

    img_rects: list = []
    for img in page.get_images(full=True):
        try:
            img_rects.extend(page.get_image_rects(img[0]))
        except Exception:  # noqa: BLE001 — xref sem rect renderizável é ignorado
            pass

    text_rects: list = []
    for b in page.get_text("blocks"):
        if len(b) > 4 and b[4].strip():
            text_rects.append(fitz.Rect(b[0], b[1], b[2], b[3]))

    return _covered_frac(img_rects), _covered_frac(text_rects)


def extract_pdf(data: bytes) -> dict:
    doc = fitz.open(stream=data, filetype="pdf")
    texts: list[str] = []
    ocr_pages: list[int] = []
    for i in range(doc.page_count):
        page = doc[i]
        text = page.get_text().strip()
        alnum = sum(c.isalnum() for c in text)
        if alnum >= MIN_PAGE_ALNUM:
            # Página com texto nativo suficiente. Ainda assim pode conter uma
            # imagem embutida com dados reais (boleto/nota renderizado como
            # imagem). Dispara OCR complementar apenas quando as imagens cobrem
            # boa parte da página e o texto nativo cobre pouco — o resultado é
            # mesclado com o texto nativo já extraído.
            texts.append(text)
            try:
                img_cov, text_cov = _page_coverage(page)
                if img_cov >= PDF_IMG_COVERAGE_MIN and text_cov <= PDF_NATIVE_TEXT_COVERAGE_MAX:
                    logger.info(
                        "OCR complementar em página híbrida: page=%d imgCov=%.2f textCov=%.2f",
                        i + 1, img_cov, text_cov,
                    )
                    supp = _ocr_page_pixmap(page)
                    if supp:
                        texts.append(supp)
                    ocr_pages.append(i + 1)
            except Exception:  # noqa: BLE001 — OCR complementar nunca regride o texto nativo
                logger.exception(
                    "falha no OCR complementar da página híbrida; mantendo só texto nativo: page=%d",
                    i + 1,
                )
        else:
            ocr_text = _ocr_page_pixmap(page)
            if ocr_text:
                texts.append(ocr_text)
            ocr_pages.append(i + 1)
    return {
        "text": "\n\n".join(t for t in texts if t),
        "pageCount": doc.page_count or 1,
        "ocrPages": ocr_pages,
    }


def extract_image(data: bytes) -> dict:
    img = cv2.imdecode(np.frombuffer(data, np.uint8), cv2.IMREAD_COLOR)
    if img is None:
        return {"text": "", "pageCount": 1, "ocrPages": []}
    return {"text": ocr_image(img), "pageCount": 1, "ocrPages": [1]}


def _extract_docx_text_via_lxml(xml_bytes: bytes) -> list[str]:
    """Extrai texto de word/document.xml usando lxml.

    O python-docx itera apenas os parágrafos e tabelas que são filhos diretos do
    body (e recursivos de tabelas), mas não entra em elementos w:sdt (Structured
    Document Tags / Content Controls). Documentos gerados pelo Word moderno — em
    especial templates com campos de formulário — armazenam praticamente todo o
    conteúdo dentro de w:sdt > w:sdtContent > w:p, tornando o resultado do
    python-docx praticamente vazio.

    Esta função usa lxml para localizar todos os w:p no documento completo (em
    qualquer profundidade, incluindo dentro de SDT e tabelas), extraindo o texto de
    cada parágrafo na ordem de documento. Isso garante que SDTs, tabelas e
    parágrafos normais sejam todos capturados.

    Retorna lista de strings não-vazias (um item por parágrafo com conteúdo).
    """
    try:
        from lxml import etree  # noqa: PLC0415 — import local para isolar dependência
    except ImportError:
        return []

    W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    try:
        tree = etree.fromstring(xml_bytes)
    except etree.XMLSyntaxError:
        return []

    parts: list[str] = []
    for p in tree.findall(f".//{{{W}}}p"):
        texts = p.findall(f".//{{{W}}}t")
        text = "".join(t.text or "" for t in texts).strip()
        if text:
            parts.append(text)
    return parts


def extract_docx(data: bytes) -> dict:
    # Extrair texto via lxml para capturar SDT (Content Controls) além de
    # parágrafos e tabelas normais. O python-docx.paragraphs ignora w:sdt,
    # causando saída vazia em templates do Word com campos de formulário.
    parts: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            if "word/document.xml" in zf.namelist():
                xml_bytes = zf.read("word/document.xml")
                parts = _extract_docx_text_via_lxml(xml_bytes)
    except (zipfile.BadZipFile, KeyError):
        pass

    # Fallback para python-docx caso lxml não esteja disponível ou falhe
    if not parts:
        doc = DocxDocument(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append("\t".join(cells))

    # OCR de imagens embutidas (word/media/*) — equivalente ao docx-images do Node.
    ocr_texts: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            for name in zf.namelist():
                if name.startswith("word/media/") and name.lower().endswith(
                    (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif")
                ):
                    img = cv2.imdecode(np.frombuffer(zf.read(name), np.uint8), cv2.IMREAD_COLOR)
                    if img is not None:
                        t = ocr_image(img)
                        if t.strip():
                            ocr_texts.append(t)
    except (zipfile.BadZipFile, KeyError):
        pass

    blocks = [b for b in ["\n".join(parts), "\n\n".join(ocr_texts)] if b.strip()]
    return {
        "text": "\n\n".join(blocks),
        "pageCount": 1,
        "ocrPages": [1] if ocr_texts else [],
    }


def extract_xlsx(data: bytes) -> dict:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    sheet_count = len(wb.worksheets)
    for ws in wb.worksheets:
        parts.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append("\t".join(cells))
    wb.close()
    return {"text": "\n".join(parts), "pageCount": sheet_count or 1, "ocrPages": []}


def extract_pptx(data: bytes) -> dict:
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    slides = list(prs.slides)
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs)
                    if txt.strip():
                        parts.append(txt)
    return {"text": "\n".join(parts), "pageCount": len(slides) or 1, "ocrPages": []}


# ──────────────────────────── endpoints ────────────────────────────


@app.get("/health")
def health() -> dict:
    return {"status": "ok", "langs": LANGS}


def _dispatch_extract(data: bytes, mime: str, name: str) -> dict:
    """Despacha a extração para o handler correto com base no MIME/nome do arquivo.
    Função síncrona; chamada via asyncio.to_thread() pelo endpoint /extract."""
    if mime == "application/pdf" or name.endswith(".pdf"):
        return extract_pdf(data)
    if mime in _IMAGE_MIMES or name.endswith(_IMAGE_EXTS):
        return extract_image(data)
    if "wordprocessingml" in mime or name.endswith(".docx"):
        return extract_docx(data)
    if "spreadsheetml" in mime or name.endswith(".xlsx"):
        return extract_xlsx(data)
    if "presentationml" in mime or name.endswith(".pptx"):
        return extract_pptx(data)
    if mime == "application/msword" or name.endswith(".doc"):
        pdf = _libreoffice_to_pdf(data, ".doc")
        return extract_pdf(pdf) if pdf else {"text": "", "pageCount": 1, "ocrPages": []}
    if mime == "application/vnd.ms-excel" or name.endswith(".xls"):
        pdf = _libreoffice_to_pdf(data, ".xls")
        return extract_pdf(pdf) if pdf else {"text": "", "pageCount": 1, "ocrPages": []}
    if mime == "application/vnd.ms-powerpoint" or name.endswith(".ppt"):
        pdf = _libreoffice_to_pdf(data, ".ppt")
        return extract_pdf(pdf) if pdf else {"text": "", "pageCount": 1, "ocrPages": []}
    if mime in ("text/rtf", "application/rtf") or name.endswith(".rtf"):
        pdf = _libreoffice_to_pdf(data, ".rtf")
        return extract_pdf(pdf) if pdf else {"text": "", "pageCount": 1, "ocrPages": []}
    if mime == "text/plain" or name.endswith(".txt"):
        return extract_txt(data)
    if mime.startswith("video/") or mime.startswith("audio/") or name.endswith(
        (".mp4", ".avi", ".mov", ".mkv", ".webm", ".3gp", ".mp3", ".m4a")
    ):
        return {"text": "", "pageCount": 1, "ocrPages": []}
    return {"text": "", "pageCount": 1, "ocrPages": [], "error": f"unsupported mime: {mime}"}


@app.post("/extract")
async def extract(
    file: UploadFile = File(...),
    content_type: str = Form(default=""),
) -> dict:
    """Extrai texto do documento recebido.

    Corre no ThreadPoolExecutor via asyncio.to_thread() para não bloquear o
    event loop do uvicorn durante OCR (Tesseract é síncrono/bloqueante).
    REQUEST_TIMEOUT_S: se o processamento ultrapassar o limite, devolve HTTP 503.
    """
    data = await file.read()
    mime = content_type or file.content_type or ""
    name = (file.filename or "").lower()

    try:
        out = await asyncio.wait_for(
            asyncio.to_thread(_dispatch_extract, data, mime, name),
            timeout=REQUEST_TIMEOUT_S,
        )
    except asyncio.TimeoutError:
        logger.error(
            "extract timeout após %ds: mime=%s name=%s bytes=%d",
            REQUEST_TIMEOUT_S, mime, name, len(data),
        )
        raise HTTPException(
            status_code=503,
            detail=(
                f"extração excedeu o tempo limite de {REQUEST_TIMEOUT_S}s — "
                "documento muito grande ou OCR lento demais neste ambiente"
            ),
        )

    out["engine"] = "python"
    return out


@app.post("/convert/pdf")
async def convert_to_pdf(
    file: UploadFile = File(...),
    content_type: str = Form(default=""),
) -> Response:
    """Converte documento Office para PDF usando LibreOffice headless."""
    data = await file.read()
    mime = content_type or file.content_type or ""
    name = (file.filename or "").lower()

    ext = OFFICE_EXTENSIONS.get(mime, "")
    if not ext:
        for suffix in OFFICE_EXTENSIONS.values():
            if name.endswith(suffix):
                ext = suffix
                break
    if not ext:
        raise HTTPException(
            status_code=422,
            detail=f"formato não suportado para conversão: {mime or name}",
        )

    pdf_bytes = _libreoffice_to_pdf(data, ext)
    if not pdf_bytes:
        raise HTTPException(status_code=500, detail="falha na conversão do documento")

    return Response(content=pdf_bytes, media_type="application/pdf")
